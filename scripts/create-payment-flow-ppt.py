#!/usr/bin/env python3
"""
토스페이먼츠 PG 심사용 결제경로 PPT 생성 스크립트
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 가맹점 정보
MERCHANT_INFO = {
    "상호명": "(주) 코로아이",
    "사업자번호": "609-81-86463",
    "URL": "https://shipedu.vercel.app",
    "Test ID": "test@shipedu.kr",
    "Test PW": "test1234!",
}

BUSINESS_INFO = {
    "상호명": "(주) 코로아이",
    "대표자명": "서중교",
    "사업자등록번호": "609-81-86463",
    "통신판매업신고번호": "(PG 계약 완료 후 신고 예정)",
    "사업장주소": "경상남도 창원시 성산구 연덕로 15번길 83(웅남동)",
    "전화번호": "055-266-8339",
}

def add_title_slide(prs, title, subtitle=""):
    """제목 슬라이드 추가"""
    slide_layout = prs.slide_layouts[6]  # 빈 레이아웃
    slide = prs.slides.add_slide(slide_layout)

    # 배경색 설정 (흰색)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.CENTER

    # 부제목
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.5))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(100, 100, 100)
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_info_slide(prs, title, info_dict, note=""):
    """정보 슬라이드 추가"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 헤더 배경
    header_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = RGBColor(0, 82, 204)  # 파란색
    header_shape.line.fill.background()

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # 정보 표시
    y_pos = 1.8
    for key, value in info_dict.items():
        # 키
        key_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(2.5), Inches(0.4))
        tf = key_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"({list(info_dict.keys()).index(key)+1}) {key}"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 82, 204)

        # 값
        val_box = slide.shapes.add_textbox(Inches(3.8), Inches(y_pos), Inches(5.5), Inches(0.4))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = f": {value}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0, 0, 0)

        y_pos += 0.5

    # 노트
    if note:
        note_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
        tf = note_box.text_frame
        p = tf.paragraphs[0]
        p.text = note
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(150, 150, 150)

    return slide

def add_screenshot_slide(prs, title, instruction, placeholder_text="[스크린샷을 여기에 붙여넣으세요]"):
    """스크린샷용 슬라이드 추가"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 헤더 배경
    header_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = RGBColor(0, 82, 204)
    header_shape.line.fill.background()

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # 설명
    inst_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(0.3))
    tf = inst_box.text_frame
    p = tf.paragraphs[0]
    p.text = instruction
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(255, 255, 255)

    # 스크린샷 자리 (점선 박스)
    placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    placeholder.fill.solid()
    placeholder.fill.fore_color.rgb = RGBColor(245, 245, 245)
    placeholder.line.color.rgb = RGBColor(200, 200, 200)
    placeholder.line.dash_style = 2  # 점선

    # 플레이스홀더 텍스트
    ph_text = slide.shapes.add_textbox(Inches(2.5), Inches(4), Inches(5), Inches(0.5))
    tf = ph_text.text_frame
    p = tf.paragraphs[0]
    p.text = placeholder_text
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.CENTER

    return slide

def create_ppt():
    """PPT 생성"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1. 표지 - 가맹점 정보
    add_title_slide(prs, "선박조종연구소", "토스페이먼츠 결제경로 파일")

    # 2. 가맹점 정보
    add_info_slide(prs, "① 가맹점 정보 기재", MERCHANT_INFO)

    # 3. 하단 정보 캡처
    add_screenshot_slide(
        prs,
        "② 하단 정보 캡처",
        "필수 구성 항목: 상호명 / 대표자명 / 사업자등록번호 / 통신판매업신고번호 / 사업장주소 / 전화번호",
        "[홈페이지 하단(Footer) 스크린샷 - 사업자 정보 포함]"
    )

    # 4. 환불규정 캡처
    add_screenshot_slide(
        prs,
        "③ 환불규정 캡처 (무형상품)",
        "환불 규정을 캡처해요. URL: https://shipedu.vercel.app/refund-policy",
        "[환불정책 페이지 스크린샷]"
    )

    # 5. 로그인/회원가입 캡처
    add_screenshot_slide(
        prs,
        "④ 로그인 / 회원가입 캡처",
        "로그인 혹은 회원가입 경로를 캡처해요. URL: https://shipedu.vercel.app/login",
        "[로그인 페이지 스크린샷]"
    )

    # 6. 상품 선택 - 강의 목록
    add_screenshot_slide(
        prs,
        "⑤ 상품 선택 / 구매과정 캡처 (1/3)",
        "강의 목록 페이지. URL: https://shipedu.vercel.app/courses",
        "[강의 목록 페이지 스크린샷]"
    )

    # 7. 상품 선택 - 강의 상세
    add_screenshot_slide(
        prs,
        "⑤ 상품 선택 / 구매과정 캡처 (2/3)",
        "강의 상세 페이지 (수강신청 버튼 포함). URL: https://shipedu.vercel.app/courses/[강의ID]",
        "[강의 상세 페이지 스크린샷 - 가격, 수강신청 버튼 표시]"
    )

    # 8. 상품 선택 - 결제 페이지
    add_screenshot_slide(
        prs,
        "⑤ 상품 선택 / 구매과정 캡처 (3/3)",
        "결제 페이지. URL: https://shipedu.vercel.app/checkout/[강의ID]",
        "[결제 페이지 스크린샷 - 구매자 정보, 결제 금액 표시]"
    )

    # 9. 카드 결제 경로 - 결제창
    add_screenshot_slide(
        prs,
        "⑥ 카드 결제경로 캡처 (1/2)",
        "토스페이먼츠 결제창 (카드 선택)",
        "[토스페이먼츠 결제창 스크린샷 - 카드 결제 선택]"
    )

    # 10. 카드 결제 경로 - 카드사 인증
    add_screenshot_slide(
        prs,
        "⑥ 카드 결제경로 캡처 (2/2)",
        "카드사 인증 화면 (비씨 카드사는 결제 직전까지 캡처 필수)",
        "[카드사 인증 화면 스크린샷]"
    )

    # 11. 마무리
    add_title_slide(prs, "감사합니다", "선박조종연구소 - (주) 코로아이")

    # 저장
    output_path = "/home/jpex/projects/private-lms-v1/docs/결제경로_선박조종연구소.pptx"
    prs.save(output_path)
    print(f"PPT 파일 생성 완료: {output_path}")
    return output_path

if __name__ == "__main__":
    create_ppt()
